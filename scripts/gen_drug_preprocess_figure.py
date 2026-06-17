"""
Drug preprocessing figure: SMILES -> 1D / 2D / 3D
Produces: drug_preprocess.pptx  (22 x 13.5 cm, compact)
Same AlphaFold/Attention-paper style as gen_arch_figure.py.

Pipeline (from scripts/pre_extract.py & src/data/dataloader.py):
  1D : AutoTokenizer -> ChemBERTa-77M-MTR last_hidden_state -> pool [354,128]
  2D : RDKit atoms/bonds -> Canonical atom/bond featurizer -> DGL bigraph
  3D : AddHs -> ETKDGv3 EmbedMultipleConfs(K) -> UFF optimize + energy -> sort
"""

import os
import math
import random
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml import etree

random.seed(7)
EMU = 360000.0
LEFT = PP_ALIGN.LEFT; CTR = PP_ALIGN.CENTER
MID = MSO_ANCHOR.MIDDLE; TOP = MSO_ANCHOR.TOP

# ── molecule image ───────────────────────────────────────────────────────────
TMP = "/tmp/tamr_assets"; os.makedirs(TMP, exist_ok=True)
MOL_PNG = os.path.join(TMP, "mol.png")
SMI = "CC(=O)Oc1ccccc1C(=O)O"
CONFS = []   # list of (png_path, energy)
try:
    from rdkit import Chem
    from rdkit.Chem import AllChem
    from rdkit.Chem.Draw import rdMolDraw2D
    # 2D depiction (input)
    m = Chem.MolFromSmiles(SMI)
    d = rdMolDraw2D.MolDraw2DCairo(420, 300)
    d.drawOptions().clearBackground = False
    rdMolDraw2D.PrepareAndDrawMolecule(d, m)
    d.FinishDrawing()
    open(MOL_PNG, "wb").write(d.GetDrawingText())
    HAVE_MOL = True
    # real 3D conformers (mirrors pre_extract.py: AddHs -> ETKDGv3 -> UFF + energy -> sort)
    import numpy as np
    from rdkit.Chem import rdMolTransforms
    m3 = Chem.AddHs(Chem.MolFromSmiles(SMI))
    p = AllChem.ETKDGv3(); p.randomSeed = 42; p.pruneRmsThresh = 0.1
    cids = list(AllChem.EmbedMultipleConfs(m3, numConfs=12, params=p))
    ens = []
    for cid in cids:
        try:
            AllChem.UFFOptimizeMolecule(m3, confId=cid, maxIters=200)
            ff = AllChem.UFFGetMoleculeForceField(m3, confId=cid)
            ens.append((cid, float(ff.CalcEnergy())))
        except Exception:
            ens.append((cid, 0.0))
    ens.sort(key=lambda t: t[1])
    disp = Chem.RemoveHs(m3)                       # heavy-atom skeleton for a cleaner 3D look
    tilts = [(0.35, 0.45), (0.95, -0.55), (-0.45, 1.05)]   # distinct 3D viewpoints
    for i, (cid, en) in enumerate(ens[:3]):
        ax, ay = tilts[i % len(tilts)]
        cx_, sx_ = math.cos(ax), math.sin(ax); cy_, sy_ = math.cos(ay), math.sin(ay)
        Rx = np.array([[1, 0, 0, 0], [0, cx_, -sx_, 0], [0, sx_, cx_, 0], [0, 0, 0, 1]])
        Ry = np.array([[cy_, 0, sy_, 0], [0, 1, 0, 0], [-sy_, 0, cy_, 0], [0, 0, 0, 1]])
        rdMolTransforms.TransformConformer(disp.GetConformer(cid), Ry.dot(Rx))
        dd = rdMolDraw2D.MolDraw2DCairo(360, 320)
        o = dd.drawOptions(); o.clearBackground = False; o.bondLineWidth = 2
        dd.DrawMolecule(disp, confId=cid)         # 3D coords (tilted) projected -> looks 3D
        dd.FinishDrawing()
        cp = os.path.join(TMP, f"conf{i}.png")
        open(cp, "wb").write(dd.GetDrawingText())
        CONFS.append((cp, en))
except Exception as e:
    print("mol gen failed:", e); HAVE_MOL = False


def rgb(r, g, b):
    return RGBColor(r, g, b)


# palette (matches arch figure)
SOFT_GREEN = rgb(216, 236, 222); GREEN_HD = rgb(38, 110, 70)
SOFT_BLUE = rgb(214, 227, 246);  BLUE_HD = rgb(40, 84, 162)
SOFT_ORANGE = rgb(250, 233, 208); ORANGE_HD = rgb(176, 104, 22)
LIGHT_BD = rgb(200, 205, 212)
TOK_GREEN = rgb(95, 180, 135); TOK_BLUE = rgb(108, 145, 216); TOK_ORANGE = rgb(238, 182, 92)
CONNECTOR = rgb(46, 84, 150)
DARK = rgb(40, 40, 40); WHITE = rgb(255, 255, 255)


def _txt(sp, text, size, bold, tcolor, align, anchor):
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_top = Cm(0.04); tf.margin_bottom = Cm(0.04)
    tf.margin_left = Cm(0.08); tf.margin_right = Cm(0.08)
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = tcolor; r.font.name = "Arial"


def box(slide, x, y, w, h, text="", size=9, bold=False, fill=WHITE, border="none",
        tcolor=DARK, align=CTR, anchor=MID, rounded=False, lw=0.9):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shp, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill == "none":
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if border == "none":
        sp.line.fill.background()
    else:
        sp.line.color.rgb = border; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    _txt(sp, text, size, bold, tcolor, align, anchor)
    return sp


def module(slide, x, y, w, h, text, fill, tcolor, size=9.5):
    return box(slide, x, y, w, h, text, size=size, bold=True, fill=fill, border="none",
               tcolor=tcolor, rounded=True)


def oval(slide, x, y, w, h, fill, border="none", lw=0.6):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(x), Cm(y), Cm(w), Cm(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if border == "none":
        sp.line.fill.background()
    else:
        sp.line.color.rgb = border; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp


def label(slide, x, y, w, text, color, size=9, align=LEFT, bold=True):
    return box(slide, x, y, w, 0.5, text, size=size, bold=bold, fill="none", border="none",
               tcolor=color, align=align, anchor=TOP)


def edge(sp, side):
    l = sp.left / EMU; t = sp.top / EMU; w = sp.width / EMU; h = sp.height / EMU
    return {'r': (l + w, t + h / 2), 'l': (l, t + h / 2),
            't': (l + w / 2, t), 'b': (l + w / 2, t + h)}[side]


def _route(s, sa, e, sb):
    x1, y1 = s; x2, y2 = e
    aligned = lambda d: abs(d) < 0.06
    if {sa, sb} <= {'r', 'l'}:
        if aligned(y1 - y2):
            return [s, e]
        xm = (x1 + x2) / 2; return [s, (xm, y1), (xm, y2), e]
    if {sa, sb} <= {'t', 'b'}:
        if aligned(x1 - x2):
            return [s, e]
        ym = (y1 + y2) / 2; return [s, (x1, ym), (x2, ym), e]
    if sa in ('r', 'l'):
        return [s, (x2, y1), e]
    return [s, (x1, y2), e]


def _round_path(pts, r=0.3, seg=8):
    if len(pts) <= 2:
        return pts
    out = [pts[0]]
    for i in range(1, len(pts) - 1):
        p0, p1, p2 = pts[i - 1], pts[i], pts[i + 1]
        v1 = (p1[0] - p0[0], p1[1] - p0[1]); v2 = (p2[0] - p1[0], p2[1] - p1[1])
        l1 = math.hypot(*v1); l2 = math.hypot(*v2)
        if l1 == 0 or l2 == 0:
            out.append(p1); continue
        rr = min(r, l1 / 2, l2 / 2)
        u1 = (v1[0] / l1, v1[1] / l1); u2 = (v2[0] / l2, v2[1] / l2)
        a = (p1[0] - u1[0] * rr, p1[1] - u1[1] * rr); b = (p1[0] + u2[0] * rr, p1[1] + u2[1] * rr)
        out.append(a)
        for k in range(1, seg):
            t = k / seg
            x = (1 - t) ** 2 * a[0] + 2 * (1 - t) * t * p1[0] + t * t * b[0]
            y = (1 - t) ** 2 * a[1] + 2 * (1 - t) * t * p1[1] + t * t * b[1]
            out.append((x, y))
        out.append(b)
    out.append(pts[-1])
    return out


def _freeform(slide, pts, color, w):
    fb = slide.shapes.build_freeform(int(round(pts[0][0] * EMU)), int(round(pts[0][1] * EMU)), scale=1)
    fb.add_line_segments([(int(round(x * EMU)), int(round(y * EMU))) for x, y in pts[1:]], close=False)
    sp = fb.convert_to_shape()
    sp.fill.background(); sp.line.color.rgb = color; sp.line.width = Pt(w); sp.shadow.inherit = False
    ln = sp._element.find('.//' + qn('a:ln'))
    ln.set('cap', 'rnd')
    etree.SubElement(ln, qn('a:round'))
    te = etree.SubElement(ln, qn('a:tailEnd'))
    te.set('type', 'arrow'); te.set('w', 'med'); te.set('len', 'med')
    return sp


def link(slide, a, sa, b, sb, color=CONNECTOR, w=1.3):
    p1 = edge(a, sa); p2 = edge(b, sb)
    return _freeform(slide, _round_path(_route(p1, sa, p2, sb)), color, w)


def token_cells(slide, x, y, n, tok, sq=0.34, gap=0.08):
    cx = x
    for _ in range(n):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(cx), Cm(y), Cm(sq), Cm(sq))
        s.fill.solid(); s.fill.fore_color.rgb = tok
        s.line.color.rgb = tok; s.line.width = Pt(0.5); s.shadow.inherit = False
        cx += sq + gap
    box(slide, cx, y - 0.06, 0.7, sq + 0.12, "...", size=11, bold=True, fill="none", border="none", align=LEFT)


def seg(slide, x1, y1, x2, y2, color, w=1.1):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    c.line.width = Pt(w); c.line.color.rgb = color; c.shadow.inherit = False


def heatmap(slide, x, y, w, h, rows, cols, base, line=WHITE):
    cw = w / cols; ch = h / rows
    for i in range(rows):
        for j in range(cols):
            v = 0.12 + 0.88 * random.random()
            shade = rgb(int(255 - (255 - base[0]) * v),
                        int(255 - (255 - base[1]) * v),
                        int(255 - (255 - base[2]) * v))
            s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Cm(x + j * cw), Cm(y + i * ch), Cm(cw), Cm(ch))
            s.fill.solid(); s.fill.fore_color.rgb = shade
            s.line.color.rgb = line; s.line.width = Pt(0.4); s.shadow.inherit = False


def graph_icon(slide, x, y, w, h):
    nodes = [(0.16, 0.34), (0.46, 0.14), (0.78, 0.30), (0.86, 0.66),
             (0.58, 0.86), (0.26, 0.74), (0.50, 0.50)]
    cols = [TOK_BLUE, TOK_BLUE, rgb(225, 120, 100), TOK_BLUE,
            rgb(225, 120, 100), TOK_BLUE, TOK_BLUE]
    pix = [(x + nx * w, y + ny * h) for nx, ny in nodes]
    edges = [(0, 1), (1, 2), (2, 3), (3, 4), (4, 5), (5, 0), (6, 1), (6, 3), (6, 5)]
    for a, b in edges:
        seg(slide, pix[a][0], pix[a][1], pix[b][0], pix[b][1], rgb(150, 168, 205), w=1.6)
    for (px, py), c in zip(pix, cols):
        oval(slide, px - 0.16, py - 0.16, 0.32, 0.32, c, border=WHITE, lw=1.0)


def adjacency_icon(slide, x, y, w, n, base):
    """sparse symmetric 0/1 adjacency; diagonal = self-loops."""
    c = w / n
    on = set()
    for i in range(n):
        for j in range(i + 1, n):
            if random.random() < 0.18:
                on.add((i, j)); on.add((j, i))
    for i in range(n):
        for j in range(n):
            filled = (i == j) or ((i, j) in on)
            shade = base if filled else rgb(236, 240, 248)
            s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x + j * c), Cm(y + i * c), Cm(c), Cm(c))
            s.fill.solid(); s.fill.fore_color.rgb = shade
            s.line.color.rgb = WHITE; s.line.width = Pt(0.4); s.shadow.inherit = False


def caption(slide, x, y, w, name, shape, note, color):
    """centered tile caption: bold name / mono shape / tiny gray note."""
    tb = box(slide, x, y, w, 1.25, "", fill="none", border="none")
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = TOP
    tf.margin_top = Cm(0.02); tf.margin_bottom = Cm(0.02); tf.margin_left = Cm(0.02); tf.margin_right = Cm(0.02)

    def addp(text, size, bold, col, font, first=False):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = CTR; p.space_after = Pt(0); p.space_before = Pt(0)
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = col; r.font.name = font

    addp(name, 7.8, True, color, "Arial", first=True)
    addp(shape, 7.8, True, DARK, "Consolas")
    if note:
        addp(note, 6.6, False, rgb(135, 135, 135), "Arial")
    return tb


def tensor_lines(slide, x, y, w, h, rows, hd):
    """key:value tensor list; rows = list of (name, shape, note)."""
    tf_box = box(slide, x, y, w, h, "", fill="none", border="none")
    tf = tf_box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MID
    tf.margin_top = Cm(0.03); tf.margin_bottom = Cm(0.03); tf.margin_left = Cm(0.05); tf.margin_right = Cm(0.05)
    for i, (name, shape, note) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = LEFT; p.space_after = Pt(3)
        r = p.add_run(); r.text = name + "  "
        r.font.size = Pt(8.5); r.font.bold = True; r.font.color.rgb = hd; r.font.name = "Arial"
        r2 = p.add_run(); r2.text = shape
        r2.font.size = Pt(8.5); r2.font.bold = True; r2.font.color.rgb = DARK; r2.font.name = "Consolas"
        if note:
            r3 = p.add_run(); r3.text = "  " + note
            r3.font.size = Pt(7.5); r3.font.bold = False; r3.font.color.rgb = rgb(120, 120, 120); r3.font.name = "Arial"
    return tf_box


# ── build ─────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Cm(23.5); prs.slide_height = Cm(15.0)
slide = prs.slides.add_slide(prs.slide_layouts[6])

L1, L2, L3 = 2.7, 7.4, 12.1   # lane center y
GREY_TX = rgb(95, 95, 95)

# SMILES input (data), feeds all three lanes
sm = box(slide, 0.6, L2 - 2.4, 4.3, 4.8, "", fill=WHITE, border=LIGHT_BD, lw=0.9)
label(slide, 0.8, L2 - 2.25, 3.0, "SMILES", DARK, size=10)
if HAVE_MOL:
    slide.shapes.add_picture(MOL_PNG, Cm(1.2), Cm(L2 - 1.6), height=Cm(2.4))
box(slide, 0.7, L2 + 1.45, 4.1, 0.8, SMI, size=8.5, bold=True,
    fill="none", border="none", tcolor=GREY_TX)

# lane spec: (center, soft, hd, tag, step1, step2)
lanes = [
    (L1, SOFT_GREEN, GREEN_HD, "① 1D · Sequence",
     "BPE Tokenizer\n(AutoTokenizer)", "ChemBERTa-77M-MTR\n(frozen)"),
    (L2, SOFT_BLUE, BLUE_HD, "② 2D · Graph",
     "RDKit parse\natoms · bonds", "Canonical atom/bond\nfeaturizer → DGL graph"),
    (L3, SOFT_ORANGE, ORANGE_HD, "③ 3D · Conformers",
     "AddHs + ETKDGv3\nembed K confs", "UFF optimize\n+ energy"),
]

s1x, s2x, ox = 5.6, 9.7, 14.0
mw, mh, ow, oh = 3.6, 1.7, 9.0, 3.9
for Lc, soft, hd, tag, t1, t2 in lanes:
    label(slide, s1x, Lc - 1.55, 7.0, tag, hd, size=9)
    m1 = module(slide, s1x, Lc - mh / 2, mw, mh, t1, soft, hd)
    m2 = module(slide, s2x, Lc - mh / 2, mw, mh, t2, soft, hd)
    out = box(slide, ox, Lc - oh / 2, ow, oh, "", fill=WHITE, border=LIGHT_BD, lw=0.8)
    link(slide, sm, 'r', m1, 'l')
    link(slide, m1, 'r', m2, 'l')
    link(slide, m2, 'r', out, 'l')

# ---- 1D output : feature-matrix heatmap ----
oy = L1 - oh / 2
label(slide, ox + 0.35, oy + 0.2, 6.0, "1D Semantic Tokens", GREEN_HD, size=9.5)
heatmap(slide, ox + 1.4, oy + 0.95, 6.2, 1.45, rows=8, cols=26, base=TOK_GREEN)
label(slide, ox + 1.4, oy + 2.55, 6.2, "Semantic tokens   [ 354 × 128 ]", GREEN_HD, size=9, align=CTR)
label(slide, ox + 1.4, oy + 2.98, 6.2, "special tokens removed · adaptive-pooled", GREY_TX, size=7, align=CTR, bold=False)

# ---- 2D output : graph + node-feats + edge-feats + adjacency ----
oy = L2 - oh / 2
label(slide, ox + 0.35, oy + 0.2, 7.0, "2D Molecular Graph   (DGL bigraph)", BLUE_HD, size=9.5)
graph_icon(slide, ox + 0.5, oy + 0.9, 1.9, 1.45)
caption(slide, ox + 0.3, oy + 2.45, 2.3, "Molecular graph", "nodes + edges", None, BLUE_HD)
# node feature matrix
heatmap(slide, ox + 2.75, oy + 0.9, 1.5, 1.45, rows=10, cols=6, base=TOK_BLUE)
caption(slide, ox + 2.5, oy + 2.45, 2.0, "Node feats  X", "[290 × 75]", None, BLUE_HD)
# edge feature matrix
heatmap(slide, ox + 4.5, oy + 0.9, 1.5, 1.45, rows=12, cols=5, base=rgb(120, 158, 224))
caption(slide, ox + 4.25, oy + 2.45, 2.0, "Edge feats  E", "[E × 12]", None, BLUE_HD)
# adjacency matrix (diagonal = self-loops)
adjacency_icon(slide, ox + 6.3, oy + 0.9, 1.45, 11, TOK_BLUE)
caption(slide, ox + 6.05, oy + 2.45, 2.0, "Adjacency  A", "[290 × 290]", None, BLUE_HD)
# feature composition (two balanced lines)
box(slide, ox + 0.3, oy + 3.26, ow - 0.6, 0.62,
    "X = 74 atom feats + 1 virtual-node bit\nE = bond type · conjugation · ring · stereo      A = self-loops on diagonal",
    size=6.8, fill="none", border="none", tcolor=GREY_TX, align=CTR, anchor=TOP)

# ---- 3D output : real RDKit conformers + energy + tensors ----
oy = L3 - oh / 2
label(slide, ox + 0.35, oy + 0.2, 7.0, "K Conformers + Energy", ORANGE_HD, size=9.5)
if CONFS:
    es = [e for _, e in CONFS]; lo = min(es); hi = max(es); rng = (hi - lo) or 1.0
    for i, (cp, en) in enumerate(CONFS):
        xi = ox + 0.45 + i * 1.5
        slide.shapes.add_picture(cp, Cm(xi), Cm(oy + 0.9), height=Cm(1.35))
        frac = 0.25 + 0.75 * ((en - lo) / rng)
        box(slide, xi + 0.1, oy + 2.4, 1.25 * frac, 0.15, "", fill=TOK_ORANGE, border="none")
        label(slide, xi + 0.1, oy + 2.55, 1.4, f"E{i+1}", ORANGE_HD, size=7.5, bold=False)
label(slide, ox + 0.45, oy + 3.05, 4.6, "K conformers  ·  sorted by energy ↑", GREY_TX, size=7.5, align=CTR, bold=False)
tensor_lines(slide, ox + 5.3, oy + 0.85, 3.6, 2.7,
             [("atom feats", "[K×64×128]", ""),
              ("coords", "[K×64×3]", "(x,y,z)"),
              ("energy", "[K]", "UFF → z-score"),
              ("conf_mask", "[K]", "valid confs")], ORANGE_HD)

print("drug preprocess figure ready")
prs.save("/Users/rifo/Code/TAMR-DTI/drug_preprocess.pptx")
