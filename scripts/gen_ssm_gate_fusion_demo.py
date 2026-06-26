"""
SSM-back-end gate-fusion demo figure.
Produces: figures/ssm_gate_fusion_demo.pptx  (22.0 x 9.5 cm)

Renders the protein-Mamba direct-gate fusion (DWFusion.ProteinMambaDirectGateBiIntention):
    P' = (1-g) * P + g * M(P)
    g  = sigmoid( b_m + MLP_g([d_bar ; p_bar ; d_bar ⊙ p_bar]) )
where M(.) is the bidirectional Mamba residual block and d_bar / p_bar are masked
max-pools of drug / protein token features.

Style helpers (palette / box / oval / freeform connectors) copied from
gen_drug_preprocess_figure.py so this script is standalone.
"""

import math
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

EMU = 360000.0
LEFT = PP_ALIGN.LEFT; CTR = PP_ALIGN.CENTER
MID = MSO_ANCHOR.MIDDLE; TOP = MSO_ANCHOR.TOP


def rgb(r, g, b): return RGBColor(r, g, b)


# palette (matches arch + drug_preprocess figures)
SOFT_BLUE = rgb(214, 227, 246);   BLUE_HD   = rgb(40, 84, 162)
SOFT_GREEN = rgb(216, 236, 222);  GREEN_HD  = rgb(38, 110, 70)
SOFT_PURPLE = rgb(229, 218, 244); PURPLE_HD = rgb(109, 68, 166)
SOFT_ORANGE = rgb(250, 233, 208); ORANGE_HD = rgb(176, 104, 22)
PANEL_FILL = rgb(248, 250, 253)
LIGHT_BD = rgb(200, 205, 212)
TOK_BLUE = rgb(108, 145, 216)
CONNECTOR = rgb(46, 84, 150)
MAMBA_LINE = rgb(115, 78, 178)
GATE_LINE = rgb(176, 104, 22)
DARK = rgb(40, 40, 40); GREY_TX = rgb(110, 110, 110); WHITE = rgb(255, 255, 255)


# ---------- primitives ----------
def _txt(sp, text, size, bold, tcolor, align, anchor, italic=False, mono=False):
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_top = Cm(0.04); tf.margin_bottom = Cm(0.04)
    tf.margin_left = Cm(0.08); tf.margin_right = Cm(0.08)
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = tcolor
        r.font.name = "Consolas" if mono else "Arial"


def box(slide, x, y, w, h, text="", size=9, bold=False, fill=WHITE, border="none",
        tcolor=DARK, align=CTR, anchor=MID, rounded=False, lw=0.9, italic=False, mono=False):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shp, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill == "none":
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if border == "none":
        sp.line.fill.background()
    else:
        sp.line.color.rgb = border; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    _txt(sp, text, size, bold, tcolor, align, anchor, italic=italic, mono=mono)
    return sp


def module(slide, x, y, w, h, text, fill, tcolor, size=10.5):
    return box(slide, x, y, w, h, text, size=size, bold=True, fill=fill,
               border="none", tcolor=tcolor, rounded=True)


def oval(slide, x, y, w, h, fill, border="none", lw=0.6, text="", tcolor=DARK,
         size=10, bold=True, italic=False):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(x), Cm(y), Cm(w), Cm(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if border == "none":
        sp.line.fill.background()
    else:
        sp.line.color.rgb = border; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    if text:
        _txt(sp, text, size, bold, tcolor, CTR, MID, italic=italic)
    return sp


def label(slide, x, y, w, text, color, size=9, align=LEFT, bold=True, mono=False):
    return box(slide, x, y, w, 0.55, text, size=size, bold=bold, fill="none",
               border="none", tcolor=color, align=align, anchor=TOP, mono=mono)


def edge(sp, side):
    l = sp.left / EMU; t = sp.top / EMU; w = sp.width / EMU; h = sp.height / EMU
    return {'r': (l + w, t + h / 2), 'l': (l, t + h / 2),
            't': (l + w / 2, t), 'b': (l + w / 2, t + h)}[side]


def _route(s, sa, e, sb):
    x1, y1 = s; x2, y2 = e
    aligned = lambda d: abs(d) < 0.06
    if {sa, sb} <= {'r', 'l'}:
        if aligned(y1 - y2):
            return [s, e]
        xm = (x1 + x2) / 2; return [s, (xm, y1), (xm, y2), e]
    if {sa, sb} <= {'t', 'b'}:
        if aligned(x1 - x2):
            return [s, e]
        ym = (y1 + y2) / 2; return [s, (x1, ym), (x2, ym), e]
    if sa in ('r', 'l'):
        return [s, (x2, y1), e]
    return [s, (x1, y2), e]


def _round_path(pts, r=0.25, seg=8):
    if len(pts) <= 2:
        return pts
    out = [pts[0]]
    for i in range(1, len(pts) - 1):
        p0, p1, p2 = pts[i - 1], pts[i], pts[i + 1]
        v1 = (p1[0] - p0[0], p1[1] - p0[1]); v2 = (p2[0] - p1[0], p2[1] - p1[1])
        l1 = math.hypot(*v1); l2 = math.hypot(*v2)
        if l1 == 0 or l2 == 0:
            out.append(p1); continue
        rr = min(r, l1 / 2, l2 / 2)
        u1 = (v1[0] / l1, v1[1] / l1); u2 = (v2[0] / l2, v2[1] / l2)
        a = (p1[0] - u1[0] * rr, p1[1] - u1[1] * rr)
        b = (p1[0] + u2[0] * rr, p1[1] + u2[1] * rr)
        out.append(a)
        for k in range(1, seg):
            t = k / seg
            x = (1 - t) ** 2 * a[0] + 2 * (1 - t) * t * p1[0] + t * t * b[0]
            y = (1 - t) ** 2 * a[1] + 2 * (1 - t) * t * p1[1] + t * t * b[1]
            out.append((x, y))
        out.append(b)
    out.append(pts[-1])
    return out


def _freeform(slide, pts, color, w, dashed=False, arrow=True):
    fb = slide.shapes.build_freeform(int(round(pts[0][0] * EMU)),
                                     int(round(pts[0][1] * EMU)), scale=1)
    fb.add_line_segments([(int(round(x * EMU)), int(round(y * EMU))) for x, y in pts[1:]],
                         close=False)
    sp = fb.convert_to_shape()
    sp.fill.background(); sp.line.color.rgb = color; sp.line.width = Pt(w)
    sp.shadow.inherit = False
    ln = sp._element.find('.//' + qn('a:ln'))
    ln.set('cap', 'rnd')
    etree.SubElement(ln, qn('a:round'))
    if dashed:
        dash = etree.SubElement(ln, qn('a:prstDash'))
        dash.set('val', 'dash')
    if arrow:
        te = etree.SubElement(ln, qn('a:tailEnd'))
        te.set('type', 'arrow'); te.set('w', 'med'); te.set('len', 'med')
    return sp


def link(slide, a, sa, b, sb, color=CONNECTOR, w=1.3, dashed=False, arrow=True):
    p1 = edge(a, sa); p2 = edge(b, sb)
    return _freeform(slide, _round_path(_route(p1, sa, p2, sb)), color, w,
                     dashed=dashed, arrow=arrow)


def token_cells(slide, x, y, n, tok, sq=0.34, gap=0.08):
    cx = x
    for _ in range(n):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(cx), Cm(y), Cm(sq), Cm(sq))
        s.fill.solid(); s.fill.fore_color.rgb = tok
        s.line.color.rgb = tok; s.line.width = Pt(0.5); s.shadow.inherit = False
        cx += sq + gap
    box(slide, cx, y - 0.06, 0.7, sq + 0.12, "...", size=11, bold=True,
        fill="none", border="none", align=LEFT)


# ── build ─────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Cm(22.0); prs.slide_height = Cm(9.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# outer panel
panel = box(slide, 0.4, 0.4, 21.2, 8.7, "", fill=PANEL_FILL, border=LIGHT_BD, lw=0.8,
            rounded=True)
label(slide, 0.85, 0.55, 18.0,
      "Post-SSM gate fusion :   Mamba-refined  M(P)   ⊕   original  P",
      DARK, size=11.5)
label(slide, 0.85, 1.10, 18.0,
      "P' = (1 - g) · P + g · M(P)        "
      "( gate g is conditioned on drug / protein context — not FiLM )",
      GREY_TX, size=9.5, bold=False, mono=True)

# ── lanes ─────────────────────────────────────────────────────────────────────
# main flow band (top): P → BiMamba → fusion → P'
TOP_Y = 2.55

# faint horizontal divider between SSM lane (top) and gate lane (bottom)
# - visual cue that drug only lives in the lower "gate" lane
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0.85), Cm(4.7),
                             Cm(20.3), Cm(0.04))
div.fill.solid(); div.fill.fore_color.rgb = LIGHT_BD
div.line.fill.background(); div.shadow.inherit = False

# left inputs: P token-strip (blue), D token-strip (green) below
P_in = box(slide, 1.0, TOP_Y, 1.4, 1.0, "P", size=15, bold=True,
           fill=SOFT_BLUE, border="none", tcolor=BLUE_HD, rounded=True)
label(slide, 0.9, TOP_Y - 0.5, 2.6, "protein tokens", GREY_TX, size=8, align=LEFT)
token_cells(slide, 1.05, TOP_Y + 1.15, 4, TOK_BLUE, sq=0.28, gap=0.06)

D_in = box(slide, 1.0, 5.95, 1.4, 1.0, "D", size=15, bold=True,
           fill=SOFT_GREEN, border="none", tcolor=GREEN_HD, rounded=True)
label(slide, 0.9, 5.45, 4.0, "drug tokens (gate lane only)",
      GREEN_HD, size=8, align=LEFT, bold=True)

# Bidirectional Mamba (purple module) - top lane
mamba = module(slide, 4.8, TOP_Y - 0.05, 4.4, 1.6,
               "Bidirectional Mamba\nM( P )    (drug-free)", SOFT_PURPLE, PURPLE_HD, size=11)

# fusion node (oval with + glyph + small caption alongside)
fuse_cx, fuse_cy = 11.6, TOP_Y + 0.75
fuse = oval(slide, fuse_cx - 0.55, fuse_cy - 0.55, 1.1, 1.1, WHITE,
            border=CONNECTOR, lw=1.6, text="+", tcolor=CONNECTOR, size=18, bold=True)
label(slide, fuse_cx + 0.7, fuse_cy - 0.5, 6.0,
      "(1 - g) · P  +  g · M(P)", DARK, size=9.5, align=LEFT, mono=True)

# refined protein P' on the right
Pp = box(slide, 17.6, TOP_Y, 3.0, 1.6,
         "Refined\nProtein  P'", size=11, bold=True,
         fill=SOFT_BLUE, border="none", tcolor=BLUE_HD, rounded=True)
token_cells(slide, 17.75, TOP_Y + 1.85, 5, TOK_BLUE, sq=0.28, gap=0.06)

# ── bottom lane: gate computation ─────────────────────────────────────────────
GATE_Y = 5.65

# masked max-pool ovals
pool_p = oval(slide, 3.0, GATE_Y - 0.55, 1.6, 1.1, WHITE, border=GATE_LINE, lw=1.4,
              text="MaxPool\np̄", tcolor=GATE_LINE, size=8.5, bold=True)
pool_d = oval(slide, 3.0, GATE_Y + 0.95, 1.6, 1.1, WHITE, border=GATE_LINE, lw=1.4,
              text="MaxPool\nd̄", tcolor=GATE_LINE, size=8.5, bold=True)

# concat node
concat = oval(slide, 5.3, GATE_Y + 0.45, 1.0, 1.0, WHITE, border=GATE_LINE, lw=1.4,
              text="‖", tcolor=GATE_LINE, size=16, bold=True)
label(slide, 5.05, GATE_Y - 0.35, 2.0, "[ d̄ ; p̄ ; d̄⊙p̄ ]", GREY_TX, size=8,
      align=CTR, mono=True, bold=False)

# gate-MLP module (orange)
gate_mlp = module(slide, 7.4, GATE_Y + 0.25, 3.0, 1.4,
                  "Gate-MLP\nLN→Lin→SiLU→Lin", SOFT_ORANGE, ORANGE_HD, size=9.5)

# add b_m (learnable bias) below MLP, small
bm = oval(slide, 11.05, GATE_Y + 0.55, 0.85, 0.85, WHITE, border=GATE_LINE, lw=1.3,
          text="b_m", tcolor=GATE_LINE, size=9, bold=True, italic=True)

# sum node
sumn = oval(slide, 12.35, GATE_Y + 0.50, 0.95, 0.95, WHITE, border=GATE_LINE, lw=1.4,
            text="+", tcolor=GATE_LINE, size=16, bold=True)

# sigmoid node (the gate value g)
sig = oval(slide, 13.85, GATE_Y + 0.50, 0.95, 0.95, WHITE, border=GATE_LINE, lw=1.6,
           text="σ", tcolor=GATE_LINE, size=16, bold=True, italic=True)
label(slide, 13.6, GATE_Y - 0.10, 1.6, "gate  g  ∈ (0,1)", GATE_LINE, size=8.5,
      align=CTR, bold=True)

# ── connectors ───────────────────────────────────────────────────────────────
# Top lane: P → Mamba (purple) → fusion (purple)
link(slide, P_in, 'r', mamba, 'l', color=MAMBA_LINE, w=1.6)
link(slide, mamba, 'r', fuse, 'l', color=MAMBA_LINE, w=1.6)

# Skip residual: P → fuse from below (navy, the (1-g)P branch)
# we route P's right edge -> down -> across -> up into fuse bottom
p_r = edge(P_in, 'r')
fuse_b = edge(fuse, 'b')
skip_pts = _round_path([
    (p_r[0] + 0.4, p_r[1]),
    (p_r[0] + 0.4, p_r[1] + 1.55),     # dip below mamba
    (fuse_b[0], p_r[1] + 1.55),
    fuse_b,
])
_freeform(slide, skip_pts, CONNECTOR, 1.4, dashed=False, arrow=True)
label(slide, p_r[0] + 0.55, p_r[1] + 1.05, 2.0, "skip  P", CONNECTOR, size=8,
      align=LEFT, bold=False, mono=True)

# Fusion → P'
link(slide, fuse, 'r', Pp, 'l', color=CONNECTOR, w=1.6)

# Bottom lane wiring
# P down to pool_p (tap from P right edge going down)
p_b = edge(P_in, 'b')
pool_p_l = edge(pool_p, 'l')
_freeform(slide, _round_path([
    p_b, (p_b[0], pool_p_l[1]), pool_p_l
]), GREY_TX, 1.1, arrow=True)

# D right to pool_d
link(slide, D_in, 'r', pool_d, 'l', color=GREY_TX, w=1.1)

# pool_p -> concat (top), pool_d -> concat (bottom)
link(slide, pool_p, 'r', concat, 'l', color=GATE_LINE, w=1.2)
link(slide, pool_d, 'r', concat, 'l', color=GATE_LINE, w=1.2)

# concat -> gate-MLP -> sum
link(slide, concat, 'r', gate_mlp, 'l', color=GATE_LINE, w=1.3)
link(slide, gate_mlp, 'r', sumn, 'l', color=GATE_LINE, w=1.3)

# b_m -> sum (from above)
link(slide, bm, 'r', sumn, 'l', color=GATE_LINE, w=1.1)

# sum -> sigmoid
link(slide, sumn, 'r', sig, 'l', color=GATE_LINE, w=1.3)

# sigmoid -> fusion (control signal, dashed orange going up into fusion top)
sig_t = edge(sig, 't')
fuse_b2 = (fuse_cx, fuse_cy + 0.55)   # use bottom of fuse
ctrl_pts = _round_path([
    sig_t,
    (sig_t[0], fuse_cy + 1.6),
    (fuse_cx, fuse_cy + 1.6),
    (fuse_cx, fuse_cy + 0.55),
])
_freeform(slide, ctrl_pts, GATE_LINE, 1.4, dashed=True, arrow=True)
label(slide, sig_t[0] + 0.05, sig_t[1] - 0.95, 1.4, "g", GATE_LINE,
      size=11, align=LEFT, bold=True, mono=True)

# ── legend strip (bottom) ─────────────────────────────────────────────────────
LEG_Y = 8.35
def legend_swatch(x, color, text):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x), Cm(LEG_Y + 0.08),
                               Cm(0.45), Cm(0.16))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    box(slide, x + 0.55, LEG_Y - 0.04, 4.5, 0.5, text, size=8.5, bold=False,
        fill="none", border="none", tcolor=GREY_TX, align=LEFT, anchor=MID)

legend_swatch(1.0,  MAMBA_LINE, "Mamba-refined branch  M(P)")
legend_swatch(7.0,  CONNECTOR,  "original-P skip branch  (and output P')")
legend_swatch(13.0, GATE_LINE,  "scalar gate g  (control signal, not data)")

# save
out_dir = "TAMR-DTI__Target-Aware_Modulation_and_Mamba-based_Representation_Refinement_for_Drug-Target_Interaction_Prediction/figures"
import os
os.makedirs(out_dir, exist_ok=True)
out = os.path.join(out_dir, "ssm_gate_fusion_demo.pptx")
prs.save(out)
print(f"wrote {out}")
