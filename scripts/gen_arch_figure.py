"""
TAMR-DTI Architecture Figure Generator (v6 — AlphaFold/Attention paper style)
Produces: tamr_dti_arch.pptx  (33.87 x 19.05 cm)

Style language (low-saturation, flat, borderless):
  - data / tensor      -> rectangle, white fill, thin light border + token cells
  - neural module      -> rounded rect, SOFT category fill, NO border, dark text
  - pointwise operator -> oval, white fill, thin border
  - section panel      -> pale tint fill, thin light-gray border, plain text title
  - connectors         -> single AF2 blue, thick, ROUND corners; modulation = red dashed
Layout stays dense (many modules) but reads clean via flat fills + aligned flows.
"""

import os
import math
import random
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml import etree

random.seed(7)
EMU = 360000.0
LEFT = PP_ALIGN.LEFT
CTR = PP_ALIGN.CENTER
MID = MSO_ANCHOR.MIDDLE
TOP = MSO_ANCHOR.TOP

# ── molecule image ───────────────────────────────────────────────────────────
TMP = "/tmp/tamr_assets"; os.makedirs(TMP, exist_ok=True)
MOL_PNG = os.path.join(TMP, "mol.png")
try:
    from rdkit import Chem
    from rdkit.Chem.Draw import rdMolDraw2D
    m = Chem.MolFromSmiles("CC(=O)Oc1ccccc1C(=O)O")
    d = rdMolDraw2D.MolDraw2DCairo(420, 300)
    d.drawOptions().clearBackground = False
    rdMolDraw2D.PrepareAndDrawMolecule(d, m)
    d.FinishDrawing()
    open(MOL_PNG, "wb").write(d.GetDrawingText())
    HAVE_MOL = True
except Exception as e:
    print("mol gen failed:", e); HAVE_MOL = False


def rgb(r, g, b):
    return RGBColor(r, g, b)


# ── palette (soft pastels for fills, darker tones only for text) ──────────────
SOFT_GREEN = rgb(216, 236, 222); GREEN_HD = rgb(38, 110, 70)
SOFT_BLUE = rgb(214, 227, 246);  BLUE_HD = rgb(40, 84, 162)
SOFT_ORANGE = rgb(250, 233, 208); ORANGE_HD = rgb(176, 104, 22)
SOFT_PURPLE = rgb(227, 220, 243); PURPLE = rgb(98, 66, 170)
SOFT_RED = rgb(247, 223, 215);   RED_HD = rgb(178, 58, 36)

PANEL_GREEN = rgb(241, 248, 244); PANEL_ORANGE = rgb(253, 247, 238)
PANEL_BLUE = rgb(239, 244, 252);  PANEL_GREY = rgb(247, 247, 249)
PANEL_RED = rgb(253, 245, 242)
PANEL_BD = rgb(206, 208, 214)
LIGHT_BD = rgb(200, 205, 212)

TOK_GREEN = rgb(95, 180, 135); TOK_ORANGE = rgb(238, 182, 92)
TOK_BLUE = rgb(108, 145, 216); TOK_RED = rgb(225, 130, 110)

CONNECTOR = rgb(46, 84, 150)
CONTRIB = rgb(186, 60, 48)
GOLD = rgb(246, 208, 84)
DARK = rgb(40, 40, 40); WHITE = rgb(255, 255, 255)


# ── primitives ───────────────────────────────────────────────────────────────
def _txt(sp, text, size, bold, tcolor, align, anchor):
    tf = sp.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_top = Cm(0.04); tf.margin_bottom = Cm(0.04)
    tf.margin_left = Cm(0.08); tf.margin_right = Cm(0.08)
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = tcolor; r.font.name = "Arial"


def box(slide, x, y, w, h, text="", size=9, bold=False, fill=WHITE, border="none",
        tcolor=DARK, align=CTR, anchor=MID, rounded=False, lw=0.9):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shp, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill == "none":
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if border == "none":
        sp.line.fill.background()
    else:
        sp.line.color.rgb = border; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    _txt(sp, text, size, bold, tcolor, align, anchor)
    return sp


def module(slide, x, y, w, h, text, fill, tcolor, size=10):
    """neural module: soft fill, no border, rounded."""
    return box(slide, x, y, w, h, text, size=size, bold=True, fill=fill,
               border="none", tcolor=tcolor, rounded=True)


def oval(slide, x, y, w, h, text="", size=8, fill=WHITE, border=LIGHT_BD, tcolor=DARK, bold=False):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(x), Cm(y), Cm(w), Cm(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = border; sp.line.width = Pt(1.0); sp.shadow.inherit = False
    _txt(sp, text, size, bold, tcolor, CTR, MID)
    return sp


def title(slide, x, y, text, color, size=11):
    return box(slide, x, y, 8.0, 0.6, text, size=size, bold=True, fill="none",
               border="none", tcolor=color, align=LEFT, anchor=TOP)


def panel(slide, x, y, w, h, fill):
    return box(slide, x, y, w, h, "", fill=fill, border=PANEL_BD, rounded=True, lw=1.0)


def edge(sp, side):
    l = sp.left / EMU; t = sp.top / EMU; w = sp.width / EMU; h = sp.height / EMU
    return {
        'r': (l + w, t + h / 2), 'l': (l, t + h / 2),
        't': (l + w / 2, t), 'b': (l + w / 2, t + h),
    }[side]


def arrow(slide, x1, y1, x2, y2, color=CONNECTOR, w=1.3, bent=False, dashed=False):
    """straight connector (used for legend / aligned flows)."""
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    c.line.width = Pt(w); c.line.color.rgb = color; c.shadow.inherit = False
    ln = c._element.find('.//' + qn('a:ln'))
    ln.set('cap', 'rnd')
    if dashed:
        etree.SubElement(ln, qn('a:prstDash')).set('val', 'dash')
    te = etree.SubElement(ln, qn('a:tailEnd'))
    te.set('type', 'arrow'); te.set('w', 'med'); te.set('len', 'med')
    return c


def _route(s, sa, e, sb):
    """orthogonal waypoints between two edge points given their exit/entry sides."""
    x1, y1 = s; x2, y2 = e
    aligned = lambda d: abs(d) < 0.06
    if {sa, sb} <= {'r', 'l'}:                         # horizontal-facing pair
        if aligned(y1 - y2):
            return [s, e]
        xm = (x1 + x2) / 2
        return [s, (xm, y1), (xm, y2), e]
    if {sa, sb} <= {'t', 'b'}:                         # vertical-facing pair
        if aligned(x1 - x2):
            return [s, e]
        ym = (y1 + y2) / 2
        return [s, (x1, ym), (x2, ym), e]
    if sa in ('r', 'l'):                               # H first, then V into top/bottom
        return [s, (x2, y1), e]
    return [s, (x1, y2), e]                            # V first, then H into left/right


def _round_path(pts, r=0.3, seg=8):
    """replace each interior corner with a smooth quadratic arc."""
    if len(pts) <= 2:
        return pts
    out = [pts[0]]
    for i in range(1, len(pts) - 1):
        p0, p1, p2 = pts[i - 1], pts[i], pts[i + 1]
        v1 = (p1[0] - p0[0], p1[1] - p0[1]); v2 = (p2[0] - p1[0], p2[1] - p1[1])
        l1 = math.hypot(*v1); l2 = math.hypot(*v2)
        if l1 == 0 or l2 == 0:
            out.append(p1); continue
        rr = min(r, l1 / 2, l2 / 2)
        u1 = (v1[0] / l1, v1[1] / l1); u2 = (v2[0] / l2, v2[1] / l2)
        a = (p1[0] - u1[0] * rr, p1[1] - u1[1] * rr)
        b = (p1[0] + u2[0] * rr, p1[1] + u2[1] * rr)
        out.append(a)
        for k in range(1, seg):
            t = k / seg
            x = (1 - t) ** 2 * a[0] + 2 * (1 - t) * t * p1[0] + t * t * b[0]
            y = (1 - t) ** 2 * a[1] + 2 * (1 - t) * t * p1[1] + t * t * b[1]
            out.append((x, y))
        out.append(b)
    out.append(pts[-1])
    return out


def _freeform(slide, pts, color, w, dashed):
    fb = slide.shapes.build_freeform(int(round(pts[0][0] * EMU)), int(round(pts[0][1] * EMU)), scale=1)
    fb.add_line_segments([(int(round(x * EMU)), int(round(y * EMU))) for x, y in pts[1:]], close=False)
    sp = fb.convert_to_shape()
    sp.fill.background()
    sp.line.color.rgb = color; sp.line.width = Pt(w); sp.shadow.inherit = False
    ln = sp._element.find('.//' + qn('a:ln'))
    ln.set('cap', 'rnd')
    if dashed:
        etree.SubElement(ln, qn('a:prstDash')).set('val', 'dash')
    etree.SubElement(ln, qn('a:round'))
    te = etree.SubElement(ln, qn('a:tailEnd'))
    te.set('type', 'arrow'); te.set('w', 'med'); te.set('len', 'med')
    return sp


def link(slide, a, sa, b, sb, color=CONNECTOR, bent=True, dashed=False, w=1.3):
    p1 = edge(a, sa); p2 = edge(b, sb)
    if not bent:
        return arrow(slide, p1[0], p1[1], p2[0], p2[1], color=color, w=w, dashed=dashed)
    pts = _round_path(_route(p1, sa, p2, sb))
    return _freeform(slide, pts, color, w, dashed)


def token_box(slide, x, y, w, h, label, n, tok, lcolor=DARK, fill=WHITE):
    """data tensor: white box, thin light border, label + colored cells."""
    sp = box(slide, x, y, w, h, "", fill=fill, border=LIGHT_BD, lw=0.8)
    if label:
        box(slide, x + 0.12, y + 0.08, w - 0.24, 0.5, label, size=8, bold=True,
            fill="none", border="none", tcolor=lcolor, align=LEFT, anchor=TOP)
    sq, gap = 0.36, 0.08
    cx = x + 0.32; sy = y + h - 0.6
    for _ in range(n):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(cx), Cm(sy), Cm(sq), Cm(sq))
        s.fill.solid(); s.fill.fore_color.rgb = tok
        s.line.color.rgb = tok; s.line.width = Pt(0.5); s.shadow.inherit = False
        cx += sq + gap
    box(slide, cx, sy - 0.06, 0.7, sq + 0.12, "...", size=11, bold=True, fill="none", border="none", align=LEFT)
    return sp


def heatmap(slide, x, y, w, h, rows=7, cols=8, base=rgb(120, 90, 200)):
    cw = w / cols; ch = h / rows
    for i in range(rows):
        for j in range(cols):
            v = random.random()
            shade = rgb(int(255 - (255 - base[0]) * v),
                        int(255 - (255 - base[1]) * v),
                        int(255 - (255 - base[2]) * v))
            s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Cm(x + j * cw), Cm(y + i * ch), Cm(cw), Cm(ch))
            s.fill.solid(); s.fill.fore_color.rgb = shade
            s.line.color.rgb = WHITE; s.line.width = Pt(0.4); s.shadow.inherit = False


def polyline(slide, pts, color=DARK, w=1.6):
    for (a, b), (c, d) in zip(pts[:-1], pts[1:]):
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(a), Cm(b), Cm(c), Cm(d))
        ln.line.width = Pt(w); ln.line.color.rgb = color; ln.shadow.inherit = False


def sigmoid_curve(slide, x, y, w, h):
    polyline(slide, [(x + 0.35, y + 0.2), (x + 0.35, y + h - 0.3), (x + w - 0.2, y + h - 0.3)],
             color=rgb(170, 170, 170), w=1.0)
    pts = []
    for k in range(25):
        t = k / 24.0
        sx = x + 0.35 + t * (w - 0.7)
        s = 1 / (1 + math.exp(-(t * 12 - 6)))
        sy = (y + h - 0.32) - s * (h - 0.62)
        pts.append((sx, sy))
    polyline(slide, pts, color=RED_HD, w=2.4)


def badge(slide, x, y, text, w=3.7):
    return box(slide, x, y, w, 0.5, text, size=8, bold=True, fill=GOLD,
               border="none", tcolor=rgb(90, 60, 0), rounded=True)


prs = Presentation()
prs.slide_width = Cm(33.87); prs.slide_height = Cm(19.05)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# ═══════════════ PANEL A · Drug Representation ═══════════════
ax, ay, aw, ah = 0.4, 0.4, 13.0, 7.3
panel(slide, ax, ay, aw, ah, PANEL_GREEN)
title(slide, ax + 0.45, ay + 0.18, "Drug Representation", GREEN_HD)

r1, r2, r3 = 2.35, 4.40, 6.45
# inputs (data)
in1 = box(slide, 0.7, r1 - 0.95, 3.5, 1.9, "", fill=WHITE, border=LIGHT_BD, lw=0.8)
box(slide, 0.85, r1 - 0.9, 2.2, 0.42, "SMILES", size=9, bold=True, fill="none",
    border="none", tcolor=GREEN_HD, align=LEFT, anchor=TOP)
if HAVE_MOL:
    slide.shapes.add_picture(MOL_PNG, Cm(1.05), Cm(r1 - 0.4), height=Cm(1.25))
in2 = box(slide, 0.7, r2 - 0.75, 3.5, 1.5, "2D Molecular\nGraph", size=10, bold=True,
          fill=WHITE, border=LIGHT_BD, tcolor=GREEN_HD, lw=0.8)
in3 = box(slide, 0.7, r3 - 0.85, 3.5, 1.7, "K Conformers\n+ Energy", size=10, bold=True,
          fill=WHITE, border=LIGHT_BD, tcolor=GREEN_HD, lw=0.8)
# encoders (module)
e1 = module(slide, 4.7, r1 - 0.75, 2.6, 1.5, "ChemBERTa", SOFT_GREEN, GREEN_HD)
e2 = module(slide, 4.7, r2 - 0.75, 2.6, 1.5, "GCN", SOFT_GREEN, GREEN_HD)
e3 = module(slide, 4.7, r3 - 0.75, 2.6, 1.5, "EGNN", SOFT_GREEN, GREEN_HD)
# token outputs (data)
t1 = token_box(slide, 7.7, r1 - 0.75, 5.1, 1.5, "Molecular Semantics", 5, TOK_GREEN, lcolor=GREEN_HD)
t2 = token_box(slide, 7.7, r2 - 0.75, 5.1, 1.5, "Topology Tokens", 5, TOK_GREEN, lcolor=GREEN_HD)
t3 = token_box(slide, 7.7, r3 - 0.75, 5.1, 1.5, "Conformer Summaries", 5, TOK_GREEN, lcolor=GREEN_HD)
for ib, eb, tb in [(in1, e1, t1), (in2, e2, t2), (in3, e3, t3)]:
    link(slide, ib, 'r', eb, 'l', bent=False)
    link(slide, eb, 'r', tb, 'l', bent=False)

# ═══════════════ PANEL B · Target-aware Conformer Weighting ═══════════════
bx, by, bw, bh = 0.4, 7.9, 13.0, 3.9
panel(slide, bx, by, bw, bh, PANEL_ORANGE)
title(slide, bx + 0.45, by + 0.16, "Target-aware Conformer Weighting", ORANGE_HD, size=10)
badge(slide, bx + bw - 3.95, by + 0.16, "★ Key Contribution ①")

pg_b = token_box(slide, 0.7, by + 0.85, 3.5, 1.1, "Protein Seed Context  p_g", 4, TOK_BLUE, lcolor=BLUE_HD)
uk_b = token_box(slide, 0.7, by + 2.2, 3.5, 1.1, "Energy Embedding  u_k", 4, TOK_ORANGE, lcolor=ORANGE_HD)
score = module(slide, 4.7, by + 1.0, 1.8, 2.0, "Score\nMLP", SOFT_ORANGE, ORANGE_HD)
box(slide, 6.9, by + 0.5, 2.6, 0.42, "α_k  (softmax)", size=9, bold=True, fill="none", border="none", tcolor=ORANGE_HD, align=LEFT)
bar_centers = [by + 1.2, by + 2.0, by + 2.8]
bars = []
for val, yc in zip([0.23, 0.47, 0.08], bar_centers):
    bw_ = 1.7 * val + 0.25
    bb = box(slide, 6.9, yc - 0.2, bw_, 0.4, "", fill=TOK_ORANGE, border="none")
    box(slide, 6.9 + bw_ + 0.1, yc - 0.22, 0.95, 0.44, f"{val:.2f}", size=8.5, fill="none", border="none", align=LEFT)
    bars.append(bb)
mul = oval(slide, 9.5, by + 1.6, 0.85, 0.85, "×", size=14, bold=True, fill=WHITE, border=ORANGE_HD, tcolor=ORANGE_HD)
geom = token_box(slide, 10.6, by + 0.85, 2.7, 2.3, "Target-aware\nGeometry Tokens", 3, TOK_ORANGE, lcolor=ORANGE_HD)
link(slide, pg_b, 'r', score, 'l', bent=True)
link(slide, uk_b, 'r', score, 'l', bent=True)
link(slide, score, 'r', bars[1], 'l', bent=False)
link(slide, bars[1], 'r', mul, 'l', bent=False)
link(slide, mul, 'r', geom, 'l', bent=False)

# ═══════════════ PANEL D · Token-wise Fusion Gate ═══════════════
dx, dy, dw, dh = 13.7, 0.4, 4.6, 11.4
panel(slide, dx, dy, dw, dh, PANEL_GREEN)
title(slide, dx + 0.35, dy + 0.16, "Token-wise", GREEN_HD, size=10)
title(slide, dx + 0.35, dy + 0.66, "Fusion Gate", GREEN_HD, size=10)
sx = dx + 0.35; sw = dw - 0.7
d_1d = token_box(slide, sx, dy + 1.6, sw, 1.1, "1D Semantics  d_1d", 4, TOK_GREEN, lcolor=GREEN_HD)
d_23 = token_box(slide, sx, dy + 3.0, sw, 1.1, "2D/3D Geometry  d_geo", 4, TOK_ORANGE, lcolor=ORANGE_HD)
gate = module(slide, sx, dy + 4.7, sw, 1.7, "", SOFT_GREEN, GREEN_HD)
box(slide, sx + 0.12, dy + 4.82, sw - 0.24, 0.4, "Sigmoid Gate (per token)", size=8, bold=True,
    fill="none", border="none", tcolor=GREEN_HD)
gxx = sx + 0.3
for v in [0.15, 0.62, 0.81, 0.33]:
    oval(slide, gxx, dy + 5.35, 0.55, 0.55, "σ", size=10, bold=True, fill=WHITE, border=GREEN_HD, tcolor=GREEN_HD)
    box(slide, gxx - 0.12, dy + 5.95, 0.8, 0.3, f"{v:.2f}", size=7, fill="none", border="none")
    gxx += 0.88
eq = oval(slide, dx + dw / 2 - 0.32, dy + 6.95, 0.64, 0.55, "=", size=13, bold=True, fill=WHITE, border=GREEN_HD)
durg = token_box(slide, sx, dy + 7.9, sw, 1.3, "DURG : Fused Drug  D", 4, TOK_GREEN, lcolor=GREEN_HD, fill=SOFT_GREEN)
lctx = token_box(slide, sx, dy + 10.0, sw, 1.3, "Ligand Context  d_g", 4, TOK_GREEN, lcolor=GREEN_HD)
link(slide, d_1d, 'b', d_23, 't', bent=False)
link(slide, d_23, 'b', gate, 't', bent=False)
link(slide, gate, 'b', eq, 't', bent=False)
link(slide, eq, 'b', durg, 't', bent=False)
link(slide, durg, 'b', lctx, 't', bent=False)

# ═══════════════ PANEL C · Protein Representation (wide bottom-left) ═══════════════
cx, cy, cw, ch = 0.4, 11.9, 17.9, 6.9
panel(slide, cx, cy, cw, ch, PANEL_BLUE)
title(slide, cx + 0.45, cy + 0.16, "Protein Representation", BLUE_HD)

seq = box(slide, 0.7, cy + 0.9, 3.0, 5.6,
          "Protein Sequence\n/ Structure\n\n[ insert BioRender\nribbon here ]\n\nMAVSEQLKVEEL\nLSKNYHLENEVAR\nLKKLV ...",
          size=8, fill=WHITE, border=LIGHT_BD, anchor=TOP, lw=0.8)
pbert = module(slide, 4.0, cy + 3.0, 2.1, 1.4, "ProtBERT", SOFT_BLUE, BLUE_HD)
p0 = token_box(slide, 6.3, cy + 1.4, 2.8, 1.4, "Protein Seed Tokens  P0", 3, TOK_BLUE, lcolor=BLUE_HD)
pg_c = token_box(slide, 6.3, cy + 4.4, 2.8, 1.4, "Protein Seed Context  p_g", 3, TOK_BLUE, lcolor=BLUE_HD)
film = module(slide, 9.4, cy + 1.4, 2.8, 4.2, "Ligand-conditioned\nFiLM\n\n\nγ  (scale)\n\nβ  (shift)", SOFT_BLUE, BLUE_HD, size=9.5)
film.text_frame.vertical_anchor = TOP
acmix = module(slide, 12.5, cy + 1.4, 2.8, 4.2, "ProteinACmix\n× 3\n\n\n( dynamic conv\n+ self-attention )", SOFT_BLUE, BLUE_HD, size=9.5)
acmix.text_frame.vertical_anchor = TOP
ptok = token_box(slide, 15.6, cy + 2.8, 2.6, 1.4, "Protein Tokens  P", 3, TOK_BLUE, lcolor=BLUE_HD, fill=SOFT_BLUE)
link(slide, seq, 'r', pbert, 'l', bent=False)
link(slide, pbert, 'r', p0, 'l', bent=True)
link(slide, p0, 'b', pg_c, 't', bent=False)
link(slide, p0, 'r', film, 'l', bent=True)
link(slide, film, 'r', acmix, 'l', bent=False)
link(slide, acmix, 'r', ptok, 'l', bent=False)

# ═══════════════ PANEL E · Protein-Mamba BiIntention Fusion ═══════════════
ex, ey, ew, eh = 18.6, 0.4, 14.9, 13.6
panel(slide, ex, ey, ew, eh, PANEL_GREY)
title(slide, ex + 0.35, ey + 0.16, "Protein-Mamba BiIntention Fusion", PURPLE)
badge(slide, ex + 10.7, ey + 0.16, "★ Key Contribution ②")

# --- protein-Mamba branch (innovation 2), single horizontal centerline ---
crow = ey + 2.3
pin = token_box(slide, ex + 0.5, crow - 0.65, 3.0, 1.3, "Protein Tokens  P", 3, TOK_BLUE, lcolor=BLUE_HD)
mamba = module(slide, ex + 4.2, crow - 1.3, 4.3, 2.6, "Bidirectional\nMamba (SSM)", SOFT_PURPLE, PURPLE, size=11)
mamba.text_frame.vertical_anchor = TOP
box(slide, ex + 4.45, crow + 0.25, 3.8, 0.36, "→  forward scan", size=7.5, fill="none", border="none", tcolor=PURPLE, align=LEFT)
box(slide, ex + 4.45, crow + 0.66, 3.8, 0.36, "←  backward scan", size=7.5, fill="none", border="none", tcolor=PURPLE, align=LEFT)
mgate = oval(slide, ex + 9.2, crow - 0.5, 1.0, 1.0, "g ⊙", size=9, bold=True, fill=WHITE, border=PURPLE, tcolor=PURPLE)
pupd = token_box(slide, ex + 10.7, crow - 0.65, 3.2, 1.3, "Refined Protein  P'", 3, TOK_BLUE, lcolor=PURPLE, fill=SOFT_PURPLE)
box(slide, ex + 4.2, ey + 3.85, 9.7, 0.42,
    "gated residual:   P  ←  P + σ(g) · ( Mamba(P) − P )", size=8.5,
    fill="none", border="none", tcolor=rgb(100, 80, 150), align=LEFT)
link(slide, pin, 'r', mamba, 'l', color=PURPLE, bent=False)
link(slide, mamba, 'r', mgate, 'l', color=PURPLE, bent=False)
link(slide, mgate, 'r', pupd, 'l', color=PURPLE, bent=False)

# --- BiIntention block (sub-container) ---
bx2, by2, bw2, bh2 = ex + 4.2, ey + 4.7, 9.6, 8.3
box(slide, bx2, by2, bw2, bh2, "", fill=PANEL_BLUE, border=PANEL_BD, rounded=True, lw=1.0)
box(slide, bx2 + 0.2, by2 + 0.14, bw2 - 0.4, 0.5, "BiIntention · Bidirectional Cross-Attention",
    size=10, bold=True, fill="none", border="none", tcolor=BLUE_HD)
arow = by2 + 1.65
d2p = module(slide, bx2 + 0.4, arow - 0.75, 4.3, 1.5, "Drug → Protein\nCross-Attn  (Q = D)", SOFT_BLUE, BLUE_HD, size=9)
p2d = module(slide, bx2 + 4.9, arow - 0.75, 4.3, 1.5, "Protein → Drug\nCross-Attn  (Q = P)", SOFT_BLUE, BLUE_HD, size=9)
din = token_box(slide, ex + 0.5, arow - 0.65, 3.0, 1.3, "Fused Drug  D", 4, TOK_GREEN, lcolor=GREEN_HD, fill=SOFT_GREEN)
hbox = box(slide, bx2 + 0.4, by2 + 3.05, 4.3, 2.6, "", fill=WHITE, border=LIGHT_BD, lw=0.8)
box(slide, bx2 + 0.55, by2 + 3.1, 4.0, 0.4, "Cross-Attention Map", size=7.8, bold=True,
    fill="none", border="none", tcolor=BLUE_HD, align=LEFT, anchor=TOP)
heatmap(slide, bx2 + 0.7, by2 + 3.62, 3.7, 1.9, rows=6, cols=10, base=rgb(108, 145, 216))
pool = module(slide, bx2 + 4.9, by2 + 3.55, 4.3, 1.6, "Max-Pool\nover tokens", SOFT_BLUE, BLUE_HD, size=9.5)
fout = token_box(slide, bx2 + 1.6, by2 + 6.4, 6.2, 1.5, "Joint Representation  f   [bs, 256]", 7, TOK_BLUE,
                 lcolor=BLUE_HD, fill=SOFT_BLUE)
link(slide, din, 'r', d2p, 'l', bent=False)
link(slide, pupd, 'b', p2d, 't', color=PURPLE, bent=True)
link(slide, d2p, 'b', hbox, 't', bent=False)
link(slide, p2d, 'b', pool, 't', bent=False)
link(slide, hbox, 'b', fout, 't', bent=True)
link(slide, pool, 'b', fout, 't', bent=True)

# --- legend (fills lower-left of E) ---
box(slide, ex + 0.45, ey + 7.7, 3.4, 5.1, "", fill=WHITE, border=LIGHT_BD, lw=0.8)
box(slide, ex + 0.6, ey + 7.85, 3.1, 0.45, "Legend", size=9.5, bold=True, fill="none", border="none", tcolor=DARK, align=LEFT, anchor=TOP)
arrow(slide, ex + 0.65, ey + 8.7, ex + 1.65, ey + 8.7, color=CONNECTOR, w=2.0)
box(slide, ex + 1.75, ey + 8.45, 2.1, 0.5, "data flow", size=8, fill="none", border="none", align=LEFT)
arrow(slide, ex + 0.65, ey + 9.6, ex + 1.65, ey + 9.6, color=CONTRIB, w=2.0, dashed=True)
box(slide, ex + 1.75, ey + 9.2, 2.1, 0.85, "target-aware\nmodulation", size=8, fill="none", border="none", tcolor=CONTRIB, align=LEFT, anchor=TOP)
oval(slide, ex + 0.85, ey + 10.45, 0.6, 0.6, "⊙", size=9, bold=True, fill=WHITE, border=PURPLE, tcolor=PURPLE)
box(slide, ex + 1.75, ey + 10.45, 2.1, 0.6, "gated fusion", size=8, fill="none", border="none", align=LEFT)
box(slide, ex + 0.75, ey + 11.35, 0.8, 0.6, "★", size=13, bold=True, fill="none", border="none", tcolor=rgb(210, 165, 30), align=CTR)
box(slide, ex + 1.75, ey + 11.35, 2.1, 0.85, "our key\ncontribution", size=8, fill="none", border="none", tcolor=rgb(150, 110, 10), align=LEFT, anchor=TOP)

# ═══════════════ PANEL F · Interaction Prediction Head ═══════════════
fx, fy, fw, fh = 18.6, 14.2, 14.9, 4.6
panel(slide, fx, fy, fw, fh, PANEL_RED)
title(slide, fx + 0.35, fy + 0.16, "Interaction Prediction Head", RED_HD)
frow = fy + 2.85
mlp = module(slide, fx + 1.6, frow - 1.5, 3.6, 3.0, "MLP Decoder\n(FC × 3 + ReLU)", SOFT_RED, RED_HD)
sig = module(slide, fx + 6.2, frow - 1.5, 3.0, 3.0, "", SOFT_RED, RED_HD)
sigmoid_curve(slide, fx + 6.45, frow - 1.2, 2.5, 2.0)
box(slide, fx + 6.2, frow + 0.92, 3.0, 0.4, "Sigmoid", size=9, bold=True, fill="none", border="none", tcolor=RED_HD)
yhat = module(slide, fx + 10.2, frow - 1.5, 4.0, 3.0, "ŷ ∈ [0, 1]\n\nInteraction\nProbability", SOFT_RED, RED_HD, size=11)
link(slide, fout, 'b', mlp, 't', bent=True)
link(slide, mlp, 'r', sig, 'l', bent=False)
link(slide, sig, 'r', yhat, 'l', bent=False)

# ═══════════════ cross-panel links ═══════════════
link(slide, t1, 'r', d_1d, 'l', bent=True)          # A -> D
link(slide, t2, 'r', d_23, 'l', bent=True)
link(slide, t3, 'b', mul, 't', bent=True)           # A -> B
link(slide, geom, 'r', d_23, 'l', bent=True)        # B -> D
link(slide, pg_c, 'l', pg_b, 'b', color=CONTRIB, bent=True, dashed=True)   # C -> B modulation
link(slide, lctx, 'l', film, 't', color=CONTRIB, bent=True, dashed=True)   # D -> C modulation
link(slide, durg, 'r', din, 'l', bent=True)         # D -> E
link(slide, ptok, 'r', pin, 'l', bent=True)         # C -> E

print("full figure ready")
prs.save("/Users/rifo/Code/TAMR-DTI/tamr_dti_arch.pptx")
